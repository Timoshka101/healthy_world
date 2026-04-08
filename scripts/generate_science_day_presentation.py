from __future__ import annotations

from pathlib import Path
from shutil import copy2

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "presentation_output"
ASSETS_DIR = OUTPUT_DIR / "assets"
TEMPLATE_PATH = Path("/home/timoshik/Downloads/Telegram Desktop/Копанев день науки.pptx")
OUTPUT_PATH = OUTPUT_DIR / "Копанев день науки - Здоровый Мир.pptx"
COPY_PATH = Path("/home/timoshik/Downloads/Telegram Desktop/Копанев день науки - Здоровый Мир.pptx")

HOME_SCREENSHOT = Path("/tmp/healthy_home.png")
GAMES_SCREENSHOT = Path("/tmp/games_auth_page.png")
DIARY_SCREENSHOT = Path("/tmp/fooddiary_page.png")
STATS_SCREENSHOT = Path("/tmp/statistics_page.png")

FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
FONT_MONO = "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf"


def font(size: int, bold: bool = False, mono: bool = False):
    path = FONT_MONO if mono else (FONT_BOLD if bold else FONT_REGULAR)
    return ImageFont.truetype(path, size)


def ensure_dirs():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def fill_rounded_rect(draw: ImageDraw.ImageDraw, xy, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def make_cover_crop(image: Image.Image, size: tuple[int, int]) -> Image.Image:
    return ImageOps.fit(image, size, method=Image.Resampling.LANCZOS, centering=(0.5, 0.0))


def make_shadow(canvas: Image.Image, box: tuple[int, int, int, int], radius: int = 28, alpha: int = 55):
    shadow = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    x1, y1, x2, y2 = box
    fill_rounded_rect(sd, (x1 + 10, y1 + 16, x2 + 10, y2 + 16), radius, (24, 40, 31, alpha))
    shadow = shadow.filter(ImageFilter.GaussianBlur(18))
    canvas.alpha_composite(shadow)


def make_browser_poster(src: Path, out: Path, title: str, crop: tuple[int, int, int, int] | None = None):
    canvas = Image.new("RGBA", (1080, 1600), (245, 247, 244, 255))
    draw = ImageDraw.Draw(canvas)

    draw.rounded_rectangle((54, 58, 1026, 1540), radius=46, fill=(255, 255, 255, 255), outline=(224, 230, 222), width=3)
    fill_rounded_rect(draw, (90, 98, 990, 188), 24, (243, 246, 241), outline=(223, 229, 221))
    for index, color in enumerate(((255, 95, 86), (255, 189, 46), (39, 201, 63))):
        draw.ellipse((120 + index * 26, 128, 136 + index * 26, 144), fill=color)
    draw.text((176, 118), title, fill=(42, 60, 52), font=font(34, bold=True))

    screenshot = Image.open(src).convert("RGB")
    if crop:
        screenshot = screenshot.crop(crop)
    screenshot = make_cover_crop(screenshot, (860, 1180))
    screen_box = (110, 220, 970, 1400)
    canvas.paste(screenshot, screen_box[:2])

    fill_rounded_rect(draw, (110, 1428, 970, 1494), 24, (247, 249, 246), outline=(226, 232, 224))
    footer = "Веб-интерфейс проекта"
    draw.text((150, 1445), footer, fill=(87, 102, 93), font=font(28))

    canvas.convert("RGB").save(out, quality=95)


def make_vitamin_poster(out: Path):
    canvas = Image.new("RGB", (1080, 1600), (245, 247, 244))
    draw = ImageDraw.Draw(canvas)

    fill_rounded_rect(draw, (60, 70, 1020, 1530), 42, (255, 255, 255), outline=(223, 229, 221), width=3)
    fill_rounded_rect(draw, (110, 120, 970, 420), 34, (142, 209, 228))
    draw.rectangle((110, 340, 970, 420), fill=(167, 217, 72))
    fill_rounded_rect(draw, (170, 300, 870, 360), 30, (72, 105, 230))
    fill_rounded_rect(draw, (170, 360, 760, 420), 30, (72, 105, 230))
    draw.text((170, 160), "Баланс витаминов", fill="white", font=font(58, bold=True))
    draw.text((170, 235), "Игрок поддерживает A, C и D в безопасной зоне", fill="white", font=font(30))

    labels = [("Vitamin A", 55, (238, 169, 58)), ("Vitamin C", 49, (225, 98, 80)), ("Vitamin D", 52, (68, 119, 223))]
    y = 540
    for label, value, color in labels:
        draw.text((130, y - 6), label, fill=(37, 50, 43), font=font(34, bold=True))
        fill_rounded_rect(draw, (360, y, 930, y + 34), 17, (232, 237, 229))
        fill_rounded_rect(draw, (360, y, 360 + int(5.7 * value), y + 34), 17, color)
        draw.text((942, y - 4), f"{value}/100", fill=(37, 50, 43), font=font(28, bold=True))
        y += 120

    draw.text((130, 920), "Продукты", fill=(37, 50, 43), font=font(38, bold=True))
    products = [("🥕", "Морковь", "A +15"), ("🍊", "Апельсин", "C +15"), ("🐟", "Рыба", "D +15"), ("🍎", "Яблоко", "C +8"), ("🥛", "Молоко", "A +8  D +8")]
    positions = [(130, 1000), (470, 1000), (130, 1210), (470, 1210), (300, 1420)]
    for (emoji, name, boost), (x, y) in zip(products, positions):
        fill_rounded_rect(draw, (x, y, x + 240, y + 150), 28, (247, 249, 246), outline=(223, 229, 221), width=3)
        draw.text((x + 26, y + 18), emoji, font=font(46), fill=(20, 20, 20))
        draw.text((x + 26, y + 76), name, font=font(30, bold=True), fill=(37, 50, 43))
        draw.text((x + 26, y + 112), boost, font=font(24), fill=(87, 102, 93))

    canvas.save(out, quality=95)


def make_tech_stack_image(out: Path):
    canvas = Image.new("RGB", (1200, 900), (250, 251, 249))
    draw = ImageDraw.Draw(canvas)
    fill_rounded_rect(draw, (50, 50, 1150, 850), 36, (255, 255, 255), outline=(226, 231, 223), width=3)
    draw.text((90, 90), "Технологический стек проекта", fill=(22, 76, 150), font=font(46, bold=True))
    draw.text((90, 155), "Backend, интерфейс, база данных и развертывание", fill=(87, 102, 93), font=font(26))

    chips = [
        ("Django 5.2", (232, 246, 236), (24, 101, 61)),
        ("PostgreSQL", (236, 242, 252), (44, 92, 173)),
        ("JavaScript", (252, 245, 214), (152, 118, 12)),
        ("Bootstrap 5", (241, 236, 251), (108, 67, 166)),
        ("HTML/CSS", (248, 239, 236), (166, 79, 45)),
        ("Docker", (233, 245, 252), (31, 125, 170)),
        ("Gunicorn", (236, 247, 238), (43, 118, 58)),
        ("Canvas API", (249, 240, 233), (149, 92, 38)),
    ]

    x_positions = [90, 360, 640]
    y = 250
    index = 0
    for row in range(3):
        for x in x_positions:
            if index >= len(chips):
                break
            label, bg, fg = chips[index]
            fill_rounded_rect(draw, (x, y, x + 220, y + 92), 28, bg)
            tw = draw.textbbox((0, 0), label, font=font(28, bold=True))[2]
            draw.text((x + (220 - tw) / 2, y + 28), label, fill=fg, font=font(28, bold=True))
            index += 1
        y += 130

    columns = [
        ("Серверная часть", ["Django views и API", "Модели и миграции", "Работа со статистикой"]),
        ("Клиентская часть", ["HTML-шаблоны", "CSS и responsive UI", "JavaScript-логика игр"]),
        ("Инфраструктура", ["Docker Compose", "Gunicorn", "PostgreSQL"]),
    ]
    start_x = 100
    for title, items in columns:
        draw.text((start_x, 640), title, fill=(37, 50, 43), font=font(30, bold=True))
        yy = 688
        for item in items:
            draw.text((start_x, yy), f"• {item}", fill=(61, 72, 66), font=font(24))
            yy += 36
        start_x += 340

    canvas.save(out, quality=95)


def make_code_image(out: Path):
    code = [
        "def calculate_macro_percentages(total_protein, total_fat, total_carb):",
        "    protein_calories = to_decimal(total_protein) * FOUR",
        "    fat_calories = to_decimal(total_fat) * NINE",
        "    carb_calories = to_decimal(total_carb) * FOUR",
        "    total_macro_calories = protein_calories + fat_calories + carb_calories",
        "",
        "    if total_macro_calories <= 0:",
        "        zero = Decimal('0.00')",
        "        return {'protein_percent': zero, 'fat_percent': zero, 'carb_percent': zero}",
        "",
        "    return {",
        "        'protein_percent': ((protein_calories / total_macro_calories) * 100).quantize(",
        "            TWO_DECIMALS, rounding=ROUND_HALF_UP",
        "        ),",
        "        'fat_percent': ((fat_calories / total_macro_calories) * 100).quantize(",
        "            TWO_DECIMALS, rounding=ROUND_HALF_UP",
        "        ),",
        "        'carb_percent': ((carb_calories / total_macro_calories) * 100).quantize(",
        "            TWO_DECIMALS, rounding=ROUND_HALF_UP",
        "        ),",
        "    }",
    ]

    canvas = Image.new("RGB", (1600, 900), (15, 21, 28))
    draw = ImageDraw.Draw(canvas)
    fill_rounded_rect(draw, (40, 40, 1560, 860), 28, (17, 24, 32))
    fill_rounded_rect(draw, (40, 40, 1560, 110), 28, (24, 35, 46))
    for index, color in enumerate(((255, 95, 86), (255, 189, 46), (39, 201, 63))):
        draw.ellipse((75 + index * 32, 64, 91 + index * 32, 80), fill=color)
    draw.text((150, 58), "games/snake_services.py", fill=(215, 222, 228), font=font(30, mono=True))

    y = 150
    for line_no, line in enumerate(code, start=1):
        draw.text((74, y), f"{line_no:>2}", fill=(94, 110, 124), font=font(24, mono=True))
        draw.text((140, y), line, fill=(231, 236, 239), font=font(24, mono=True))
        y += 33

    canvas.save(out, quality=95)


def create_assets():
    make_browser_poster(HOME_SCREENSHOT, ASSETS_DIR / "home_card.png", "Главная страница", crop=(20, 0, 1340, 840))
    make_browser_poster(DIARY_SCREENSHOT, ASSETS_DIR / "diary_card.png", "Дневник питания", crop=(20, 0, 1340, 1260))
    make_browser_poster(GAMES_SCREENSHOT, ASSETS_DIR / "games_card.png", "Игровой раздел", crop=(20, 0, 1340, 1280))
    make_browser_poster(STATS_SCREENSHOT, ASSETS_DIR / "stats_card.png", "Статистика игр", crop=(20, 0, 1340, 1180))
    make_browser_poster(GAMES_SCREENSHOT, ASSETS_DIR / "snake_card.png", "Змейка про здоровое питание", crop=(40, 320, 1320, 1210))
    make_vitamin_poster(ASSETS_DIR / "vitamin_card.png")
    make_tech_stack_image(ASSETS_DIR / "tech_stack.png")
    make_code_image(ASSETS_DIR / "code_macro.png")


def set_text(shape, text: str, max_size: int = 32):
    shape.text = text
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    try:
        text_frame.fit_text(font_family="Arial", max_size=max_size)
    except Exception:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(max_size)


def largest_pictures(slide, count: int):
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    return sorted(pictures, key=lambda s: s.width * s.height, reverse=True)[:count]


def replace_picture(slide, shape, image_path: Path):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_presentation():
    prs = Presentation(TEMPLATE_PATH)

    # Slide 1
    set_text(prs.slides[0].shapes[0], "Разработка веб-приложения\n«Здоровый Мир»", 34)

    # Slide 2
    set_text(prs.slides[1].shapes[0], "Методологические характеристики исследования", 28)
    set_text(
        prs.slides[1].shapes[1],
        "Объект исследования – процессы цифрового сопровождения питания, учета рациона и формирования здоровых пищевых привычек пользователей.",
        26,
    )

    # Slide 3
    set_text(prs.slides[2].shapes[0], "Методологические характеристики исследования", 28)
    set_text(
        prs.slides[2].shapes[1],
        "Предмет исследования – веб-технологии, модели данных и игровые механики, применяемые для контроля калорийности, БЖУ и обучения основам здорового питания.",
        25,
    )

    # Slide 4
    set_text(prs.slides[3].shapes[0], "Методологические характеристики исследования", 28)
    set_text(
        prs.slides[3].shapes[1],
        "Цель исследования – разработать веб-приложение «Здоровый Мир» для ведения дневника питания, расчета целевых показателей и изучения принципов здорового питания в интерактивной форме.",
        24,
    )

    # Slide 5
    set_text(prs.slides[4].shapes[2], "Задачи исследования", 28)
    set_text(
        prs.slides[4].shapes[0],
        "1. Проанализировать предметную область здорового питания и пользовательские сценарии.\n"
        "2. Спроектировать структуру данных для продуктов, дневника, игровых уровней и результатов.\n"
        "3. Реализовать веб-интерфейс с авторизацией, дневником питания и настройками пользователя.\n"
        "4. Разработать обучающие мини-игры и модуль статистики результатов.\n"
        "5. Подготовить проект к запуску в Docker-окружении и демонстрации на сервере.",
        21,
    )

    # Slide 6
    set_text(prs.slides[5].shapes[0], "Структура веб-приложения", 30)
    set_text(
        prs.slides[5].shapes[1],
        "МОДУЛЬ 1. Пользовательская часть\n"
        "1.1 Регистрация, вход и профиль пользователя\n"
        "1.2 Настройки целей питания и персональные данные\n\n"
        "МОДУЛЬ 2. Функциональная часть\n"
        "2.1 Дневник питания и расчет калорий, белков, жиров и углеводов\n"
        "2.2 Игровой раздел, статистика прохождения и серверное API",
        20,
    )

    # Slide 7
    set_text(prs.slides[6].shapes[2], "Разработка ПО", 28)
    set_text(
        prs.slides[6].shapes[1],
        "Для разработки веб-приложения были выбраны современные инструменты:\n"
        "– Django 5.2 как основной backend-фреймворк;\n"
        "– PostgreSQL для хранения данных пользователей и результатов;\n"
        "– HTML, CSS, Bootstrap 5 и JavaScript для интерфейса;\n"
        "– Docker и Docker Compose для развёртывания;\n"
        "– Gunicorn для запуска приложения на сервере.",
        22,
    )
    for picture in largest_pictures(prs.slides[6], 1):
        replace_picture(prs.slides[6], picture, ASSETS_DIR / "tech_stack.png")

    # Slide 8
    set_text(prs.slides[7].shapes[0], "Разработка ПО", 28)
    set_text(
        prs.slides[7].shapes[1],
        "В ходе разработки было создано веб-приложение «Здоровый Мир». "
        "Оно объединяет учёт питания, расчёт целевых норм и интерактивные обучающие игры.",
        24,
    )
    set_text(
        prs.slides[7].shapes[2],
        "Основные компоненты\n"
        "— пользовательская авторизация и профиль;\n"
        "— дневник питания с поиском и добавлением продуктов;\n"
        "— расчёт суточной нормы калорий и БЖУ;\n"
        "— игровой раздел со «Змейкой», «Балансом витаминов» и тестами;\n"
        "— статистика результатов и история прохождений.",
        22,
    )

    # Slide 9
    set_text(prs.slides[8].shapes[0], "Разработка ПО", 28)
    set_text(prs.slides[8].shapes[1], "Функционал приложения", 22)
    set_text(
        prs.slides[8].shapes[3],
        "Функционал\n"
        "— регистрация и вход пользователя;\n"
        "— ведение дневника питания по дням;\n"
        "— расчёт калорийности и макронутриентов по каждому приёму пищи;\n"
        "— настройка целей в зависимости от пола, веса, роста и активности;\n"
        "— игра «Змейка про здоровое питание» с подсчётом БЖУ;\n"
        "— игра «Баланс витаминов» и два тестовых модуля;\n"
        "— сохранение истории попыток и общей игровой статистики.",
        20,
    )

    # Slide 10
    set_text(prs.slides[9].shapes[0], "Разработка ПО", 28)
    set_text(prs.slides[9].shapes[1], "Пример ключевой бизнес-логики", 22)
    set_text(
        prs.slides[9].shapes[2],
        "На слайде представлен фрагмент серверной логики, вычисляющей процентное соотношение "
        "белков, жиров и углеводов. Этот расчёт используется в игровых сценариях, при оценке "
        "результата и в аналитике пользователя.",
        20,
    )
    for picture in largest_pictures(prs.slides[9], 1):
        replace_picture(prs.slides[9], picture, ASSETS_DIR / "code_macro.png")

    # Slide 11
    set_text(prs.slides[10].shapes[0], "Разработка интерфейса", 28)
    set_text(
        prs.slides[10].shapes[6],
        "На слайде представлены основные экраны веб-приложения: главная страница, дневник питания "
        "и игровой раздел. Интерфейс разработан с упором на наглядность, простую навигацию и "
        "быстрый доступ к основным функциям.",
        20,
    )
    slide11_pics = largest_pictures(prs.slides[10], 3)
    for shape, asset in zip(slide11_pics, [ASSETS_DIR / "home_card.png", ASSETS_DIR / "diary_card.png", ASSETS_DIR / "games_card.png"]):
        replace_picture(prs.slides[10], shape, asset)

    # Slide 12
    set_text(prs.slides[11].shapes[0], "Игровой модуль и статистика", 28)
    slide12_pics = largest_pictures(prs.slides[11], 3)
    for shape, asset in zip(slide12_pics, [ASSETS_DIR / "snake_card.png", ASSETS_DIR / "vitamin_card.png", ASSETS_DIR / "stats_card.png"]):
        replace_picture(prs.slides[11], shape, asset)

    # Slide 13
    set_text(prs.slides[12].shapes[0], "Заключение", 28)
    set_text(
        prs.slides[12].shapes[1],
        "В ходе выполнения работы было разработано веб-приложение «Здоровый Мир», объединяющее "
        "дневник питания, расчёт персональных норм и обучающие мини-игры. "
        "В приложении реализованы хранение данных в PostgreSQL, игровой API, интерфейс для "
        "пользователя и модуль статистики. Проект готов к демонстрации и дальнейшему расширению.",
        22,
    )

    # Slide 14
    set_text(prs.slides[13].shapes[0], "Спасибо за внимание!", 34)

    prs.save(OUTPUT_PATH)
    copy2(OUTPUT_PATH, COPY_PATH)


if __name__ == "__main__":
    ensure_dirs()
    from PIL import ImageFilter  # imported lazily to keep startup tidy

    create_assets()
    build_presentation()
    print(f"Saved presentation to: {OUTPUT_PATH}")
    print(f"Copied presentation to: {COPY_PATH}")
